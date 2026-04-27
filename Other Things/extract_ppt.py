import collections
import collections.abc
import pptx

def extract_text(filename):
    print(f"--- Extracting text from {filename} ---")
    try:
        prs = pptx.Presentation(filename)
        for i, slide in enumerate(prs.slides):
            print(f"\n[Slide {i+1}]")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
    except Exception as e:
        print(f"Error reading {filename}: {e}")

extract_text("STTG_Presentation.pptx")
extract_text("PPT Format for RRP.pptx")
